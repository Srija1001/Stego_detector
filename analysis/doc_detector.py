"""
Document & Media Steganography Detector
=========================================
Covers: PDF, DOCX, PPTX, XLSX, TXT, CSV, HTML, XML, MD, RTF, ODT, ODS, ODP
        Audio: MP3, WAV, FLAC, OGG, M4A, AAC, AIFF
        Video: MP4, AVI, MKV, MOV, WMV, FLV, WEBM

Detection techniques
--------------------
Plain text (TXT/CSV/MD/HTML/XML/RTF):
  - Zero-width invisible Unicode characters
  - Trailing-space whitespace encoding
  - Unicode homoglyph substitution
  - Mixed / anomalous line endings
  - Null bytes

PDF:
  - Embedded files / JavaScript / hidden actions in raw stream
  - Invisible / white-colour text layers  (needs PyMuPDF)
  - Metadata unusual keys

DOCX / ODT:
  - Hidden text (w:vanish attribute)
  - White-on-white text runs
  - VBA macros
  - Custom XML data parts
  - Zero-width characters in XML

PPTX / ODP:
  - Hidden slides (show=0)
  - Off-canvas / out-of-bounds shapes
  - White text boxes
  - Speaker notes hidden characters
  - Zero-width characters in XML

XLSX / ODS:
  - Hidden sheets (hidden / veryHidden)
  - Hidden rows / columns
  - White-font or tiny-font cells (< 4 pt)
  - VBA macros
  - Zero-width characters in XML

Audio (WAV/MP3/FLAC/OGG/M4A):
  - LSB encoding in raw audio samples  (WAV/FLAC)
  - Excess metadata / ID3 tags
  - Statistical randomness in sample LSBs
  - Appended data after audio end marker
  - Unusually high noise floor in silent segments

Video (MP4/AVI/MKV/MOV etc.):
  - Extra data appended after container end
  - Multiple audio/video streams (hidden stream)
  - Unusually large metadata blocks
  - File size vs duration anomaly
  - Suspicious embedded attachments (MKV)
"""

import os
import zipfile
import hashlib
import struct
import wave
from pathlib import Path

# ── optional imports ─────────────────────────────────────────────────────────
try:
    import fitz
    HAS_PYMUPDF = True
except ImportError:
    HAS_PYMUPDF = False

try:
    from docx import Document as DocxDoc
    from docx.oxml.ns import qn
    HAS_DOCX = True
except ImportError:
    HAS_DOCX = False

try:
    from pptx import Presentation
    HAS_PPTX = True
except ImportError:
    HAS_PPTX = False

try:
    import openpyxl
    HAS_OPENPYXL = True
except ImportError:
    HAS_OPENPYXL = False

try:
    import numpy as np
    import librosa
    import soundfile as sf
    HAS_AUDIO = True
except ImportError:
    HAS_AUDIO = False

try:
    import cv2
    HAS_CV2 = True
except ImportError:
    HAS_CV2 = False

# ── Zero-width Unicode characters ────────────────────────────────────────────
ZERO_WIDTH = {
    '\u200b': 'ZERO WIDTH SPACE',
    '\u200c': 'ZERO WIDTH NON-JOINER',
    '\u200d': 'ZERO WIDTH JOINER',
    '\u200e': 'LTR MARK',
    '\u200f': 'RTL MARK',
    '\u2060': 'WORD JOINER',
    '\u2062': 'INVISIBLE TIMES',
    '\u2063': 'INVISIBLE SEPARATOR',
    '\ufeff': 'BOM / ZERO WIDTH NO-BREAK SPACE',
    '\u00ad': 'SOFT HYPHEN',
    '\u202a': 'LTR EMBEDDING',
    '\u202b': 'RTL EMBEDDING',
    '\u202c': 'POP DIRECTIONAL FORMATTING',
    '\u202d': 'LTR OVERRIDE',
    '\u202e': 'RTL OVERRIDE',
}

# Cyrillic / lookalike homoglyphs for Latin letters
HOMOGLYPHS = {
    '\u0430': 'a', '\u0435': 'e', '\u043e': 'o', '\u0440': 'p',
    '\u0441': 'c', '\u0443': 'y', '\u0445': 'x',
    '\u0456': 'i', '\u0458': 'j',
}

SUPPORTED_EXTENSIONS = {
    # documents
    '.pdf', '.docx', '.doc', '.pptx', '.ppt',
    '.xlsx', '.xls', '.txt', '.csv', '.html',
    '.htm', '.xml', '.rtf', '.md', '.odt', '.odp', '.ods',
    # audio
    '.wav', '.mp3', '.flac', '.ogg', '.m4a', '.aac', '.aiff',
    # video
    '.mp4', '.avi', '.mkv', '.mov', '.wmv', '.flv', '.webm',
}

IMAGE_EXTENSIONS = {'.png', '.jpg', '.jpeg', '.bmp', '.tiff', '.tif', '.webp', '.gif'}

AUDIO_EXTENSIONS = {'.wav', '.mp3', '.flac', '.ogg', '.m4a', '.aac', '.aiff'}
VIDEO_EXTENSIONS = {'.mp4', '.avi', '.mkv', '.mov', '.wmv', '.flv', '.webm'}
PLAIN_TEXT_EXTS  = {'.txt', '.csv', '.md', '.html', '.htm', '.xml', '.rtf'}


# ── Shared helpers ────────────────────────────────────────────────────────────

def _sha256(path):
    h = hashlib.sha256()
    with open(path, 'rb') as f:
        for chunk in iter(lambda: f.read(65536), b''):
            h.update(chunk)
    return h.hexdigest()

def _count_zw(text):
    return {name: text.count(ch) for ch, name in ZERO_WIDTH.items() if text.count(ch)}

def _decode_zw_binary(text, max_chars=500):
    """Attempt to decode zero-width character binary steganography.
    ZWSP (U+200B)=0, ZWNJ (U+200C)=1, ZWJ (U+200D) may also encode bits."""
    # Build bit stream from ZW chars
    ZW_BITS = {'\u200b': '0', '\u200c': '1', '\u200d': '1'}
    bits = ''.join(ZW_BITS.get(c, '') for c in text)
    if len(bits) < 8:
        return ''
    chars = []
    for i in range(0, len(bits) - 7, 8):
        byte_val = int(bits[i:i+8], 2)
        if byte_val == 0:
            break
        if 32 <= byte_val <= 126:
            chars.append(chr(byte_val))
        else:
            chars.append('.')
        if len(chars) >= max_chars:
            break
    result = ''.join(chars)
    printable = sum(1 for c in result if c != '.') / max(len(result), 1)
    return result if printable > 0.5 else ''

def _decode_homoglyphs(text):
    """Extract the message encoded via homoglyph substitution."""
    decoded = []
    for c in text:
        if c in HOMOGLYPHS:
            decoded.append(f'[{HOMOGLYPHS[c]}]')
        elif c.isprintable():
            decoded.append(c)
    return ''.join(decoded)

def _count_homoglyphs(text):
    return {f'{ch}→{lat}': text.count(ch) for ch, lat in HOMOGLYPHS.items() if text.count(ch)}

def _decode_lsb_samples_text(int_array, max_chars=500):
    """Decode a null-terminated ASCII message from the LSBs of an integer
    sample array (used for WAV PCM audio LSB steganography)."""
    bits = (int_array & 1).astype(np.uint8)
    chars, null_found = [], False
    for i in range(0, len(bits) - 7, 8):
        bv = int(np.packbits(bits[i:i+8])[0])
        if bv == 0:
            null_found = True
            break
        chars.append(chr(bv) if 32 <= bv <= 126 else '.')
        if len(chars) >= max_chars:
            break
    text = ''.join(chars)
    pr = sum(1 for c in text if c != '.') / max(1, len(text))
    return {
        'text': text, 'length': len(text),
        'printable_ratio': round(pr, 3),
        'null_terminated': null_found,
        'likely_message': pr > 0.7 and len(text) > 4,
    }

def _bytes_text_preview(raw_bytes, max_chars=500):
    """Render a printable preview of raw appended/chunk bytes, replacing
    non-printable bytes with '.', for display as potential hidden content."""
    if not raw_bytes:
        return {'text': '', 'printable_ratio': 0.0, 'likely_message': False}
    sample = raw_bytes[:max_chars]
    text = ''.join(chr(b) if 32 <= b <= 126 else '.' for b in sample)
    pr = sum(1 for c in text if c != '.') / max(1, len(text))
    return {
        'text': text, 'length': len(text),
        'printable_ratio': round(pr, 3),
        'likely_message': pr > 0.6 and len(text.strip('.')) > 4,
    }

def _trailing_ws(text):
    lines = text.splitlines()
    t = sum(1 for l in lines if l != l.rstrip(' '))
    return {'total_lines': len(lines), 'trailing': t,
            'ratio': round(t / max(len(lines), 1), 4)}

def _read_text(path):
    raw = open(path, 'rb').read()
    for enc in ('utf-8', 'utf-16', 'latin-1'):
        try:
            return raw.decode(enc, errors='replace')
        except Exception:
            pass
    return raw.decode('latin-1', errors='replace')

def _zip_xml(path):
    """Return all XML text from an Office Open XML zip."""
    out = ''
    try:
        with zipfile.ZipFile(path, 'r') as z:
            for name in z.namelist():
                if name.endswith('.xml') or name.endswith('.rels'):
                    out += z.read(name).decode('utf-8', errors='replace')
    except Exception:
        pass
    return out

def _risk_level(score):
    if score >= 70: return 'HIGH'
    if score >= 35: return 'MEDIUM'
    if score >= 10: return 'LOW'
    return 'CLEAN'

def _assemble(path, analyzer_label, checks):
    """Build the standard result dict from a checks dict."""
    total = min(100, sum(c.get('score', 0) for c in checks.values()))
    detected_list = [k for k, v in checks.items() if v.get('detected')]
    return {
        'analyzer_type':   analyzer_label,
        'file_path':       path,
        'file_type':       Path(path).suffix.lower(),
        'file_size_bytes': os.path.getsize(path),
        'file_hash':       _sha256(path),
        'checks':          checks,
        'total_score':     total,
        'risk_level':      _risk_level(total),
        'detected':        total >= 35,
        'detected_checks': detected_list,
        'summary': (
            f'{len(detected_list)} indicator(s): {", ".join(detected_list)}'
            if detected_list else 'No steganography indicators found.'
        ),
    }


# ─────────────────────────────────────────────────────────────────────────────
#  Plain Text
# ─────────────────────────────────────────────────────────────────────────────
class PlainTextAnalyzer:
    def __init__(self, path):
        self.path = path
        self._text = _read_text(path)
        self._raw  = open(path, 'rb').read()

    def run_all(self):
        zw      = _count_zw(self._text)
        zw_n    = sum(zw.values())
        hg      = _count_homoglyphs(self._text)
        hg_n    = sum(hg.values())
        tws     = _trailing_ws(self._text)
        tws_det = tws['ratio'] > 0.30 or tws['trailing'] > 50
        crlf    = self._raw.count(b'\r\n')
        lf      = self._raw.count(b'\n') - crlf
        cr      = self._raw.count(b'\r') - crlf
        mixed   = (crlf > 0 and lf > 0) or cr > 0
        nulls   = self._raw.count(b'\x00')

        checks = {
            'zero_width_chars': {
                'found': zw, 'total': zw_n,
                'detected': bool(zw),
                'decoded_content': _decode_zw_binary(self._text) if zw_n > 8 else '',
                'score': min(65, zw_n * 8),
                'info': f'{zw_n} invisible Unicode character(s)',
            },
            'trailing_whitespace': {
                'stats': tws, 'detected': tws_det,
                'score': min(40, int(tws['ratio'] * 100)) if tws_det else 0,
                'info': f'{tws["trailing"]} lines with trailing spaces ({tws["ratio"]:.1%})',
            },
            'homoglyphs': {
                'found': hg, 'total': hg_n,
                'detected': bool(hg),
                'score': min(55, hg_n * 8),
                'info': f'{hg_n} lookalike character(s)',
            },
            'line_ending_anomaly': {
                'CRLF': crlf, 'LF': lf, 'CR': cr, 'mixed': mixed,
                'detected': mixed,
                'score': 15 if mixed else 0,
                'info': 'Mixed line endings detected' if mixed else 'Consistent line endings',
            },
            'null_bytes': {
                'count': nulls, 'detected': bool(nulls),
                'score': min(70, nulls * 10),
                'info': f'{nulls} null byte(s)',
            },
        }
        return _assemble(self.path, 'Plain Text', checks)


# ─────────────────────────────────────────────────────────────────────────────
#  PDF
# ─────────────────────────────────────────────────────────────────────────────
class PDFAnalyzer:
    def __init__(self, path):
        self.path = path
        self._raw = open(path, 'rb').read()

    def run_all(self):
        # Raw stream scan
        suspicious = {}
        raw_lo = self._raw.lower()
        for name, patterns in {
            'javascript':     [b'/javascript', b'/js '],
            'embedded_files': [b'/embeddedfile', b'/filespec'],
            'launch_action':  [b'/launch'],
            'openaction':     [b'/openaction'],
        }.items():
            n = sum(raw_lo.count(p) for p in patterns)
            if n:
                suspicious[name] = n

        raw_text = self._raw.decode('latin-1', errors='replace')
        zw = _count_zw(raw_text)
        zw_n = sum(zw.values())
        zw_decoded = _decode_zw_binary(raw_text) if zw_n > 8 else ''

        stream_score = min(30, len(suspicious) * 8) + (min(30, zw_n * 3) if zw else 0)
        checks = {
            'raw_stream_scan': {
                'suspicious_keys': suspicious, 'zero_width_chars': zw,
                'detected': bool(suspicious) or bool(zw),
                'decoded_content': zw_decoded,
                'score': stream_score,
                'info': (f'Suspicious PDF keys: {list(suspicious.keys())}; ZW chars: {zw_n}'
                         if suspicious or zw else 'No suspicious stream content'),
            },
        }

        # Invisible text (PyMuPDF)
        if HAS_PYMUPDF:
            try:
                doc = fitz.open(self.path)
                invis = 0
                invis_texts = []
                for page in doc:
                    for block in page.get_text('dict').get('blocks', []):
                        for line in block.get('lines', []):
                            for span in line.get('spans', []):
                                if span.get('color') == 16777215 or span.get('size', 12) < 0.5:
                                    txt = span.get('text', '')
                                    invis += len(txt)
                                    if txt.strip():
                                        invis_texts.append(txt)
                total_pages = doc.page_count
                blank = sum(1 for p in doc if len(p.get_text().strip()) < 5)
                doc.close()
                checks['invisible_text'] = {
                    'invisible_chars': invis, 'detected': invis > 0,
                    'decoded_content': ' '.join(invis_texts).strip(),
                    'score': min(65, invis // 3),
                    'info': f'{invis} invisible/white character(s)',
                }
                checks['blank_pages'] = {
                    'total': total_pages, 'blank': blank,
                    'detected': blank > 0 and blank / max(total_pages, 1) > 0.5,
                    'score': min(20, blank * 5),
                    'info': f'{blank}/{total_pages} blank pages',
                }
            except Exception as e:
                checks['invisible_text'] = {'detected': False, 'score': 0, 'info': f'Error: {e}'}
        else:
            checks['invisible_text'] = {
                'detected': False, 'score': 0,
                'info': 'PyMuPDF not installed (pip install pymupdf)',
            }

        return _assemble(self.path, 'PDF', checks)


# ─────────────────────────────────────────────────────────────────────────────
#  DOCX / ODT
# ─────────────────────────────────────────────────────────────────────────────
class DOCXAnalyzer:
    def __init__(self, path):
        self.path = path

    def run_all(self):
        checks = {}

        # Hidden text (vanish)
        if HAS_DOCX:
            try:
                doc = DocxDoc(self.path)
                hidden, samples, full_texts = 0, [], []
                for para in doc.paragraphs:
                    for run in para.runs:
                        rpr = run._r.find(qn('w:rPr'))
                        if rpr is not None and rpr.find(qn('w:vanish')) is not None:
                            hidden += 1
                            if run.text.strip():
                                full_texts.append(run.text)
                            if len(samples) < 10:
                                samples.append(run.text[:120])
                full_content = ' '.join(full_texts)
                checks['hidden_text'] = {
                    'count': hidden, 'samples': samples,
                    'full_content': full_content,
                    'detected': hidden > 0,
                    'score': min(70, hidden * 15),
                    'info': f'{hidden} hidden (vanish) text run(s)',
                }
            except Exception as e:
                checks['hidden_text'] = {'detected': False, 'score': 0, 'info': f'Error: {e}'}

            # White text
            try:
                doc = DocxDoc(self.path)
                white = 0
                white_texts = []
                for para in doc.paragraphs:
                    for run in para.runs:
                        try:
                            rgb = run.font.color.rgb
                            if rgb and str(rgb).upper() in ('FFFFFF', 'FFFFFE'):
                                white += 1
                                if run.text.strip():
                                    white_texts.append(run.text)
                        except Exception:
                            pass
                checks['white_text'] = {
                    'count': white, 'detected': white > 0,
                    'full_content': ' '.join(white_texts),
                    'score': min(60, white * 12),
                    'info': f'{white} white-coloured text run(s)',
                }
            except Exception as e:
                checks['white_text'] = {'detected': False, 'score': 0, 'info': f'Error: {e}'}
        else:
            checks['hidden_text'] = {'detected': False, 'score': 0, 'info': 'python-docx not installed'}
            checks['white_text']  = {'detected': False, 'score': 0, 'info': 'python-docx not installed'}

        # Zero-width chars in XML
        xml_text = _zip_xml(self.path)
        zw = _count_zw(xml_text)
        zw_n = sum(zw.values())
        zw_decoded = _decode_zw_binary(xml_text) if zw_n > 8 else ''
        checks['zero_width_chars'] = {
            'found': zw, 'total': zw_n, 'detected': bool(zw),
            'decoded_content': zw_decoded,
            'score': min(65, zw_n * 8),
            'info': f'{zw_n} zero-width character(s) in document XML',
        }

        # VBA macros
        try:
            with zipfile.ZipFile(self.path, 'r') as z:
                has_macro = any('vbaProject' in n for n in z.namelist())
            checks['vba_macros'] = {
                'detected': has_macro,
                'score': 40 if has_macro else 0,
                'info': 'VBA macro project present' if has_macro else 'No macros',
            }
        except Exception as e:
            checks['vba_macros'] = {'detected': False, 'score': 0, 'info': f'Error: {e}'}

        # Custom XML parts
        try:
            with zipfile.ZipFile(self.path, 'r') as z:
                parts = [n for n in z.namelist() if 'customXml' in n]
            checks['custom_xml'] = {
                'parts': parts, 'detected': bool(parts),
                'score': min(30, len(parts) * 10),
                'info': f'{len(parts)} custom XML part(s)',
            }
        except Exception as e:
            checks['custom_xml'] = {'detected': False, 'score': 0, 'info': f'Error: {e}'}

        return _assemble(self.path, 'DOCX/ODT', checks)


# ─────────────────────────────────────────────────────────────────────────────
#  PPTX / ODP
# ─────────────────────────────────────────────────────────────────────────────
class PPTXAnalyzer:
    def __init__(self, path):
        self.path = path

    def run_all(self):
        checks = {}

        if HAS_PPTX:
            try:
                prs = Presentation(self.path)

                # Hidden slides
                hidden_slides = [i for i, sl in enumerate(prs.slides)
                                 if sl._element.get('show') == '0']
                hidden_slide_texts = []
                for i in hidden_slides:
                    slide = prs.slides[i]
                    for shape in slide.shapes:
                        if shape.has_text_frame and shape.text_frame.text.strip():
                            hidden_slide_texts.append(shape.text_frame.text)
                checks['hidden_slides'] = {
                    'indices': hidden_slides, 'count': len(hidden_slides),
                    'detected': bool(hidden_slides),
                    'decoded_content': '\n---\n'.join(hidden_slide_texts),
                    'score': min(65, len(hidden_slides) * 20),
                    'info': f'{len(hidden_slides)} hidden slide(s)',
                }

                # Off-canvas shapes
                sw, sh = prs.slide_width, prs.slide_height
                offcanvas = 0
                offcanvas_texts = []
                for slide in prs.slides:
                    for shape in slide.shapes:
                        l = shape.left or 0
                        t = shape.top  or 0
                        w = shape.width or 0
                        h = shape.height or 0
                        if l + w < 0 or l > sw or t + h < 0 or t > sh:
                            offcanvas += 1
                            if shape.has_text_frame and shape.text_frame.text.strip():
                                offcanvas_texts.append(shape.text_frame.text)
                checks['offcanvas_shapes'] = {
                    'count': offcanvas, 'detected': offcanvas > 0,
                    'decoded_content': '\n---\n'.join(offcanvas_texts),
                    'score': min(55, offcanvas * 15),
                    'info': f'{offcanvas} shape(s) outside slide boundary',
                }

                # White text
                white = 0
                white_texts = []
                for slide in prs.slides:
                    for shape in slide.shapes:
                        if shape.has_text_frame:
                            for para in shape.text_frame.paragraphs:
                                for run in para.runs:
                                    try:
                                        from pptx.dml.color import RGBColor
                                        if run.font.color.rgb == RGBColor(0xFF, 0xFF, 0xFF):
                                            white += 1
                                            if run.text.strip():
                                                white_texts.append(run.text)
                                    except Exception:
                                        pass
                checks['white_text'] = {
                    'count': white, 'detected': white > 0,
                    'decoded_content': ' '.join(white_texts),
                    'score': min(55, white * 12),
                    'info': f'{white} white-text run(s)',
                }

                # Notes steganography
                notes_text = ''
                for slide in prs.slides:
                    if slide.has_notes_slide:
                        notes_text += slide.notes_slide.notes_text_frame.text
                zw_n  = _count_zw(notes_text)
                hg_n  = _count_homoglyphs(notes_text)
                notes_total = sum(zw_n.values()) + sum(hg_n.values())
                notes_decoded = _decode_zw_binary(notes_text) if sum(zw_n.values()) > 8 else ''
                checks['notes_stego'] = {
                    'zero_width': zw_n, 'homoglyphs': hg_n,
                    'detected': bool(zw_n) or bool(hg_n),
                    'decoded_content': notes_decoded,
                    'score': min(50, notes_total * 5),
                    'info': f'{notes_total} hidden character(s) in speaker notes',
                }
            except Exception as e:
                for k in ('hidden_slides', 'offcanvas_shapes', 'white_text', 'notes_stego'):
                    checks[k] = {'detected': False, 'score': 0, 'info': f'Error: {e}'}
        else:
            for k in ('hidden_slides', 'offcanvas_shapes', 'white_text', 'notes_stego'):
                checks[k] = {'detected': False, 'score': 0, 'info': 'python-pptx not installed'}

        # Zero-width chars in XML
        xml_text = _zip_xml(self.path)
        zw = _count_zw(xml_text)
        zw_n = sum(zw.values())
        checks['zero_width_chars'] = {
            'found': zw, 'total': zw_n, 'detected': bool(zw),
            'score': min(65, zw_n * 8),
            'info': f'{zw_n} zero-width character(s) in XML',
        }

        return _assemble(self.path, 'PPTX/ODP', checks)


# ─────────────────────────────────────────────────────────────────────────────
#  XLSX / ODS
# ─────────────────────────────────────────────────────────────────────────────
class XLSXAnalyzer:
    def __init__(self, path):
        self.path = path

    def run_all(self):
        checks = {}

        if HAS_OPENPYXL:
            try:
                wb = openpyxl.load_workbook(self.path, read_only=False, data_only=True)
                hidden_sheets = [s for s in wb.sheetnames
                                 if wb[s].sheet_state in ('hidden', 'veryHidden')]
                hidden_content = []
                for sname in hidden_sheets:
                    ws = wb[sname]
                    for row in ws.iter_rows():
                        for cell in row:
                            if cell.value not in (None, ''):
                                hidden_content.append(str(cell.value))
                wb.close()
                checks['hidden_sheets'] = {
                    'sheets': hidden_sheets, 'detected': bool(hidden_sheets),
                    'decoded_content': '\n'.join(hidden_content[:200]),
                    'score': min(65, len(hidden_sheets) * 20),
                    'info': f'{len(hidden_sheets)} hidden sheet(s): {hidden_sheets}',
                }
            except Exception as e:
                checks['hidden_sheets'] = {'detected': False, 'score': 0, 'info': f'Error: {e}'}

            try:
                wb = openpyxl.load_workbook(self.path)
                hr = hc = white = tiny = 0
                white_texts = []
                hidden_rc_texts = []
                for ws in wb.worksheets:
                    hidden_row_idx = {idx for idx, rd in ws.row_dimensions.items() if rd.hidden}
                    hidden_col_idx = {idx for idx, cd in ws.column_dimensions.items() if cd.hidden}
                    hr += len(hidden_row_idx)
                    hc += len(hidden_col_idx)
                    for row in ws.iter_rows():
                        for cell in row:
                            if cell.value not in (None, ''):
                                font = cell.font
                                if (cell.row in hidden_row_idx or
                                        (cell.column_letter in hidden_col_idx)):
                                    hidden_rc_texts.append(str(cell.value))
                                if font:
                                    if font.color and font.color.rgb:
                                        if str(font.color.rgb).upper() in ('FFFFFFFF', '00FFFFFF'):
                                            white += 1
                                            white_texts.append(str(cell.value))
                                    if font.size and font.size < 4:
                                        tiny += 1
                wb.close()
                checks['hidden_rows_cols'] = {
                    'hidden_rows': hr, 'hidden_cols': hc,
                    'detected': hr > 0 or hc > 0,
                    'decoded_content': '\n'.join(hidden_rc_texts[:200]),
                    'score': min(50, (hr + hc) * 4),
                    'info': f'{hr} hidden row(s), {hc} hidden column(s)',
                }
                checks['white_or_tiny_cells'] = {
                    'white_cells': white, 'tiny_font_cells': tiny,
                    'decoded_content': '\n'.join(white_texts[:200]),
                    'detected': white > 0 or tiny > 0,
                    'score': min(60, (white + tiny) * 10),
                    'info': f'{white} white-font cell(s), {tiny} tiny-font cell(s)',
                }
            except Exception as e:
                checks['hidden_rows_cols']    = {'detected': False, 'score': 0, 'info': f'Error: {e}'}
                checks['white_or_tiny_cells'] = {'detected': False, 'score': 0, 'info': f'Error: {e}'}
        else:
            for k in ('hidden_sheets', 'hidden_rows_cols', 'white_or_tiny_cells'):
                checks[k] = {'detected': False, 'score': 0, 'info': 'openpyxl not installed'}

        # Zero-width chars
        xml_text = _zip_xml(self.path)
        zw = _count_zw(xml_text)
        zw_n = sum(zw.values())
        checks['zero_width_chars'] = {
            'found': zw, 'total': zw_n, 'detected': bool(zw),
            'score': min(65, zw_n * 8),
            'info': f'{zw_n} zero-width character(s)',
        }

        # VBA macros
        try:
            with zipfile.ZipFile(self.path, 'r') as z:
                has_macro = any('vbaProject' in n for n in z.namelist())
            checks['vba_macros'] = {
                'detected': has_macro,
                'score': 40 if has_macro else 0,
                'info': 'VBA macro project present' if has_macro else 'No macros',
            }
        except Exception as e:
            checks['vba_macros'] = {'detected': False, 'score': 0, 'info': f'Error: {e}'}

        return _assemble(self.path, 'XLSX/ODS', checks)


# ─────────────────────────────────────────────────────────────────────────────
#  Audio
# ─────────────────────────────────────────────────────────────────────────────
class AudioAnalyzer:
    """
    Audio steganography detection.

    Techniques detected:
    1. LSB randomness  — natural audio has correlated LSBs; stego destroys this.
    2. Excess metadata — unusually large ID3/INFO chunks can hide data.
    3. Appended data   — bytes after the audio end marker.
    4. Silent segment noise — hidden data raises the noise floor in silent parts.
    """

    def __init__(self, path):
        self.path = path
        self.ext  = Path(path).suffix.lower()

    def run_all(self):
        checks = {}

        if not HAS_AUDIO:
            # FIX 1: Don't bail out entirely — still run numpy-based WAV LSB check
            # and all structural checks below even without librosa.
            pass

        # 1. Waveform smoothness + LSB analysis
        # FIX 1: Added numpy-based WAV PCM chi-square LSB check that works WITHOUT
        # librosa. The original code returned detected=False whenever librosa was
        # missing, silently skipping WAV LSB detection entirely.
        try:
            if self.ext == '.wav':
                # Direct PCM LSB chi-square — no librosa needed
                import wave as _wave
                with _wave.open(self.path, 'rb') as _w:
                    _frames = _w.readframes(_w.getnframes())
                    _sampwidth = _w.getsampwidth()
                # 16-bit samples are most common; fall back to raw bytes for others
                if _sampwidth == 2:
                    _samples = np.frombuffer(_frames, dtype=np.int16)
                else:
                    _samples = np.frombuffer(_frames, dtype=np.uint8)
                lsb_bits = _samples & 1
                ones_ratio = float(np.mean(lsb_bits))
                # Chi-square: natural audio LSBs are NOT perfectly 50/50
                # Stego (LSB substitution) forces them very close to 0.5000
                # Chi-square on LSB pairs: natural audio LSBs correlate with neighbours
                # Stego LSBs are independent message bits → near-zero autocorrelation
                lsb_autocorr = float(np.corrcoef(lsb_bits[:9999].astype(float),
                                                  lsb_bits[1:10000].astype(float))[0, 1])
                lsb_suspicious = abs(lsb_autocorr) < 0.008 or abs(ones_ratio - 0.5) < 0.002
                # Also check adjacent-sample correlation on raw PCM
                x16 = _samples[:50000].astype(np.float32)
                sample_corr = float(np.corrcoef(x16[:-1], x16[1:])[0, 1]) if len(x16) > 1 else 1.0
                corr_suspicious = sample_corr < 0.5
                lsb_sus    = lsb_suspicious  # alias for info string below
                suspicious = lsb_suspicious or corr_suspicious

                # Attempt to decode an actual hidden message from the LSBs
                _decoded = _decode_lsb_samples_text(_samples[:2000 * 8 + 16], max_chars=2000) \
                    if suspicious else {'text': '', 'likely_message': False}

                checks['lsb_randomness'] = {
                    'ones_ratio':        round(ones_ratio, 5),
                    'sample_correlation': round(sample_corr, 4),
                    'detected':          suspicious,
                    'score':             60 if suspicious else 0,
                    'decoded_content':   _decoded.get('text', ''),
                    'decoded_likely_message': _decoded.get('likely_message', False),
                    'decoded_printable_ratio': _decoded.get('printable_ratio', 0),
                    'decoded_null_terminated': _decoded.get('null_terminated', False),
                    'info': (
                        f'LSB ones_ratio={ones_ratio:.5f}, lsb_autocorr={lsb_autocorr:.5f} '
                        f'sample corr={sample_corr:.3f} '
                        f'— {"SUSPICIOUS: uniform LSBs indicate LSB embedding" if lsb_suspicious else ""}'
                        f'{"noise-like waveform" if corr_suspicious and not lsb_suspicious else ""}'
                        f'{"normal audio waveform" if not suspicious else ""}'
                    ),
                }
                checks['silent_segment_noise'] = {'detected': False, 'score': 0, 'info': 'N/A for WAV PCM path'}

            elif HAS_AUDIO:
                samples, sr = librosa.load(self.path, sr=None, mono=True)
                n = min(len(samples), 50000)
                x = samples[:n]

                if len(x) > 1:
                    sample_corr = float(np.corrcoef(x[:-1], x[1:])[0, 1])
                else:
                    sample_corr = 1.0

                diff_energy  = float(np.mean(np.diff(x) ** 2))
                total_energy = float(np.mean(x ** 2)) + 1e-10
                diff_ratio   = diff_energy / total_energy
                suspicious   = sample_corr < 0.5 and diff_ratio > 1.0

                checks['lsb_randomness'] = {
                    'sample_correlation': round(sample_corr, 4),
                    'diff_energy_ratio':  round(diff_ratio, 4),
                    'detected':           suspicious,
                    'score':              50 if suspicious else 0,
                    'info': (
                        f'Sample correlation={sample_corr:.3f}, diff-energy ratio={diff_ratio:.3f} '
                        f'— {"noise-like waveform (possible LSB stego)" if suspicious else "normal audio waveform"}'
                    ),
                }

                # 4. Silent segment noise floor
                try:
                    threshold   = 0.01 * float(np.max(np.abs(samples)) + 1e-10)
                    silent_mask = np.abs(samples) < threshold
                    if silent_mask.sum() > sr * 0.1:
                        silent_rms = float(np.sqrt(np.mean(samples[silent_mask] ** 2)))
                        noise_suspicious = silent_rms > 0.005
                        checks['silent_segment_noise'] = {
                            'silent_samples': int(silent_mask.sum()),
                            'rms_in_silence': round(silent_rms, 6),
                            'detected':       noise_suspicious,
                            'score':          30 if noise_suspicious else 0,
                            'info': (
                                f'RMS in silent segments: {silent_rms:.6f} '
                                f'— {"elevated noise floor" if noise_suspicious else "normal silence level"}'
                            ),
                        }
                    else:
                        checks['silent_segment_noise'] = {
                            'detected': False, 'score': 0,
                            'info': 'Not enough silence to analyse',
                        }
                except Exception as e:
                    checks['silent_segment_noise'] = {'detected': False, 'score': 0, 'info': f'Error: {e}'}
            else:
                checks['lsb_randomness'] = {'detected': False, 'score': 0, 'info': 'librosa not installed and not a WAV file'}
                checks['silent_segment_noise'] = {'detected': False, 'score': 0, 'info': 'Skipped'}

        except Exception as e:
            checks['lsb_randomness'] = {'detected': False, 'score': 0, 'info': f'Load error: {e}'}
            checks['silent_segment_noise'] = {'detected': False, 'score': 0, 'info': 'Skipped'}

        # 2. Metadata / header size
        try:
            raw = open(self.path, 'rb').read(65536)  # first 64 KB
            file_size = os.path.getsize(self.path)

            meta_bytes = 0
            if self.ext == '.mp3':
                # ID3 tag size is encoded in bytes 6-9 (synchsafe int)
                if raw[:3] == b'ID3' and len(raw) >= 10:
                    b6, b7, b8, b9 = raw[6], raw[7], raw[8], raw[9]
                    meta_bytes = (b6 << 21) | (b7 << 14) | (b8 << 7) | b9
            elif self.ext == '.wav':
                # Check INFO chunk size
                if raw[:4] == b'RIFF' and len(raw) >= 12:
                    meta_bytes = raw.count(b'LIST') * 100  # rough estimate

            meta_ratio = meta_bytes / max(file_size, 1)
            meta_suspicious = meta_bytes > 50000 or meta_ratio > 0.3
            checks['excess_metadata'] = {
                'metadata_bytes': meta_bytes,
                'file_size':      file_size,
                'ratio':          round(meta_ratio, 4),
                'detected':       meta_suspicious,
                'score':          40 if meta_suspicious else 0,
                'info': (
                    f'Metadata: {meta_bytes} bytes ({meta_ratio:.1%} of file)'
                    f' — {"unusually large" if meta_suspicious else "normal"}'
                ),
            }
        except Exception as e:
            checks['excess_metadata'] = {'detected': False, 'score': 0, 'info': f'Error: {e}'}

        # 3. Appended data after audio end
        try:
            raw_full  = open(self.path, 'rb').read()
            file_size = len(raw_full)
            appended  = 0
            marker    = None

            if self.ext == '.mp3':
                # ID3v1 tag is last 128 bytes if it starts with 'TAG'
                if raw_full[-128:-125] == b'TAG':
                    end_pos = file_size - 128
                else:
                    end_pos = file_size
                # Find last valid MP3 sync word (0xFF 0xE0 mask)
                for i in range(file_size - 4, max(file_size - 2000, 0), -1):
                    if raw_full[i] == 0xFF and (raw_full[i+1] & 0xE0) == 0xE0:
                        appended = end_pos - i - 4
                        break
            elif self.ext == '.wav':
                # WAV: RIFF chunk size is bytes 4-7 (little-endian)
                if raw_full[:4] == b'RIFF' and file_size >= 8:
                    declared = struct.unpack_from('<I', raw_full, 4)[0] + 8
                    appended = max(0, file_size - declared)
                    marker = 'RIFF size field'
            elif self.ext == '.flac':
                # FIX 5: Implement FLAC appended data detection via sync word search.
                # Original hardcoded appended=0. FLAC frames start with 0xFFF8/0xFFF9.
                last_sync = raw_full.rfind(b'\xff\xf8')
                if last_sync == -1:
                    last_sync = raw_full.rfind(b'\xff\xf9')
                # 128-byte buffer accounts for the last frame's own bytes
                appended = max(0, file_size - last_sync - 128) if last_sync != -1 else 0

            elif self.ext in ('.ogg',):
                # FIX 6: OGG appended data detection (was never implemented).
                # Find last OGG capture pattern 'OggS' and calculate page end.
                last_ogg = raw_full.rfind(b'OggS')
                if last_ogg != -1 and last_ogg + 27 < file_size:
                    seg_count  = raw_full[last_ogg + 26]
                    seg_table  = raw_full[last_ogg + 27: last_ogg + 27 + seg_count]
                    page_end   = last_ogg + 27 + seg_count + sum(seg_table)
                    appended   = max(0, file_size - page_end)
                else:
                    appended = 0

            appended = max(0, appended)
            # FIX: Lowered threshold from 512 to 16 bytes — even tiny appended data is notable
            suspicious = appended > 16
            _appended_preview = (
                _bytes_text_preview(raw_full[file_size - appended:], max_chars=1000)
                if suspicious and appended > 0 else {'text': '', 'likely_message': False}
            )
            checks['appended_data'] = {
                'appended_bytes': appended,
                'detected': suspicious,
                'score': min(70, appended // 10) if suspicious else 0,
                'decoded_content': _appended_preview.get('text', ''),
                'decoded_likely_message': _appended_preview.get('likely_message', False),
                'info': (
                    f'{appended} bytes after audio end'
                    f' — {"suspicious hidden data" if suspicious else "normal"}'
                ),
            }
        except Exception as e:
            checks['appended_data'] = {'detected': False, 'score': 0, 'info': f'Error: {e}'}

        # FIX 2: WAV unknown chunk scanner.
        # Data injected as a non-standard chunk INSIDE the RIFF container has
        # declared size == actual size, so appended_data always returns 0.
        # This scanner walks every chunk and flags non-standard types.
        if self.ext == '.wav':
            try:
                raw_w = open(self.path, 'rb').read()
                KNOWN_WAV_CHUNKS = {
                    b'fmt ', b'data', b'LIST', b'INFO', b'fact',
                    b'cue ', b'JUNK', b'bext', b'plst', b'ltxt',
                    b'smpl', b'inst', b'id3 ', b'ID3 ', b'_PMX',
                }
                pos, unknown_chunks = 12, []
                while pos + 8 <= len(raw_w):
                    ctype = raw_w[pos:pos + 4]
                    if len(ctype) < 4: break
                    csize = struct.unpack_from('<I', raw_w, pos + 4)[0]
                    if ctype not in KNOWN_WAV_CHUNKS:
                        content   = raw_w[pos + 8: pos + 8 + min(csize, 2000)]
                        printable = sum(1 for b in content if 32 <= b <= 126) / max(len(content), 1)
                        unknown_chunks.append({
                            'type': ctype.decode('ascii', 'replace'),
                            'size': csize, 'offset': pos,
                            'printable_ratio': round(printable, 3),
                            'content_bytes': content,
                        })
                    pos += 8 + csize + (csize % 2)
                    if pos > len(raw_w): break
                suspicious_chunks = [c for c in unknown_chunks if c['printable_ratio'] > 0.2 or c['size'] > 32]
                chunk_suspicious  = len(suspicious_chunks) > 0
                _chunk_decoded = ''
                if suspicious_chunks:
                    best_chunk = max(suspicious_chunks, key=lambda c: c['printable_ratio'])
                    _preview = _bytes_text_preview(best_chunk['content_bytes'], max_chars=1000)
                    _chunk_decoded = _preview.get('text', '')
                # Drop raw bytes from the stored chunk list (keep JSON exports lean)
                for c in unknown_chunks:
                    c.pop('content_bytes', None)
                checks['unknown_chunks'] = {
                    'detected': chunk_suspicious,
                    'score':    65 if chunk_suspicious else 0,
                    'chunks':   unknown_chunks,
                    'decoded_content': _chunk_decoded,
                    'decoded_likely_message': bool(_chunk_decoded.strip('.')),
                    'info': (
                        f'Unknown WAV chunks: {[c["type"] for c in unknown_chunks]} '
                        f'— {"SUSPICIOUS: non-standard chunks may hide data" if chunk_suspicious else "all standard"}'
                        if unknown_chunks else 'All standard WAV chunks'
                    ),
                }
            except Exception as e:
                checks['unknown_chunks'] = {'detected': False, 'score': 0, 'info': f'Error: {e}'}

        return _assemble(self.path, 'Audio', checks)


# ─────────────────────────────────────────────────────────────────────────────
#  Video
# ─────────────────────────────────────────────────────────────────────────────
class VideoAnalyzer:
    """
    Video steganography detection.

    Techniques:
    1. Appended data after container end (easy hiding spot).
    2. Multiple streams — extra hidden audio/video track.
    3. File size vs duration anomaly — too large for its duration.
    4. Large metadata blocks.
    5. MKV attachments (legitimate containers sometimes, but can hide files).
    """

    def __init__(self, path):
        self.path = path
        self.ext  = Path(path).suffix.lower()

    def run_all(self):
        checks = {}

        # 1. Appended data
        try:
            raw_full  = open(self.path, 'rb').read()
            file_size = len(raw_full)
            appended  = 0

            if self.ext in ('.mp4', '.mov'):
                # FIX 3a: Correct MP4 atom end — size counts from atom START (pos-4),
                # not from the 4-byte type field (pos). Old code used pos+atom_size
                # which skips the 4-byte size field itself, under-counting by 4 bytes.
                # Also removed the arbitrary 100-byte tolerance that hid small payloads.
                # Check both moov and mdat; whichever ends later is the true container end.
                container_end = 0
                for marker in (b'moov', b'mdat'):
                    mpos = raw_full.rfind(marker)
                    if mpos != -1 and mpos >= 4:
                        atom_size = struct.unpack_from('>I', raw_full, mpos - 4)[0]
                        end = (mpos - 4) + atom_size   # atom_start + full atom size
                        container_end = max(container_end, end)
                appended = max(0, file_size - container_end) if container_end else 0
            elif self.ext == '.avi':
                # AVI: RIFF container
                if raw_full[:4] == b'RIFF' and file_size >= 8:
                    declared = struct.unpack_from('<I', raw_full, 4)[0] + 8
                    appended = max(0, file_size - declared)

            # FIX 3b: Lowered threshold from 1024 to 16 bytes
            suspicious = appended > 16
            _v_preview = (
                _bytes_text_preview(raw_full[file_size - appended:], max_chars=1000)
                if suspicious and appended > 0 else {'text': '', 'likely_message': False}
            )
            checks['appended_data'] = {
                'appended_bytes': appended,
                'detected': suspicious,
                'score': min(65, appended // 500) if suspicious else 0,
                'decoded_content': _v_preview.get('text', ''),
                'decoded_likely_message': _v_preview.get('likely_message', False),
                'info': (
                    f'{appended} bytes after container end'
                    f' — {"possible hidden data" if suspicious else "normal"}'
                ),
            }
        except Exception as e:
            checks['appended_data'] = {'detected': False, 'score': 0, 'info': f'Error: {e}'}

        # 2. Stream count & 3. Size/duration anomaly (OpenCV)
        if HAS_CV2:
            try:
                cap = cv2.VideoCapture(self.path)
                fps      = cap.get(cv2.CAP_PROP_FPS) or 1
                frames   = cap.get(cv2.CAP_PROP_FRAME_COUNT) or 0
                w        = int(cap.get(cv2.CAP_PROP_FRAME_WIDTH))
                h        = int(cap.get(cv2.CAP_PROP_FRAME_HEIGHT))
                cap.release()

                duration_s  = frames / fps if fps > 0 else 0
                file_size   = os.path.getsize(self.path)

                # Expected size: rough estimate (uncompressed would be huge; compressed is usually
                # < 50 MB/minute for HD). Flag if > 200 MB/min
                expected_max = duration_s * 200 * 1024 * 1024 / 60
                size_suspicious = duration_s > 1 and file_size > expected_max

                checks['size_duration_anomaly'] = {
                    'duration_s':   round(duration_s, 2),
                    'file_size_mb': round(file_size / (1024*1024), 2),
                    'fps': round(fps, 2),
                    'resolution': f'{w}x{h}',
                    'detected': size_suspicious,
                    'score': 35 if size_suspicious else 0,
                    'info': (
                        f'{duration_s:.1f}s video, {file_size/(1024*1024):.1f} MB'
                        f' — {"unusually large for duration" if size_suspicious else "normal ratio"}'
                    ),
                }
            except Exception as e:
                checks['size_duration_anomaly'] = {'detected': False, 'score': 0, 'info': f'Error: {e}'}
        else:
            checks['size_duration_anomaly'] = {
                'detected': False, 'score': 0,
                'info': 'OpenCV not installed (pip install opencv-python)',
            }

        # 4. Large metadata blocks
        try:
            raw_start = open(self.path, 'rb').read(65536)
            file_size = os.path.getsize(self.path)

            # Count metadata-like markers
            meta_markers = {
                b'ftyp': 'MP4 file type',
                b'udta': 'MP4 user data',
                b'meta': 'metadata atom',
                b'XMP_': 'XMP metadata',
                b'ID3':  'ID3 tag',
            }
            found_meta = {name: raw_start.count(marker)
                          for marker, name in meta_markers.items()
                          if raw_start.count(marker) > 0}

            # Simple heuristic: if meta atoms take > 10% of early data, flag
            meta_bytes_est = sum(v * 500 for v in found_meta.values())
            meta_suspicious = meta_bytes_est > 5000 or ('XMP_' in str(found_meta) and found_meta.get('XMP_', 0) > 2)
            checks['large_metadata'] = {
                'found_markers': found_meta,
                'detected': meta_suspicious,
                'score': 25 if meta_suspicious else 0,
                'info': (
                    f'Metadata markers: {found_meta}'
                    f' — {"possibly excessive" if meta_suspicious else "normal"}'
                ),
            }
        except Exception as e:
            checks['large_metadata'] = {'detected': False, 'score': 0, 'info': f'Error: {e}'}

        # FIX 4: AVI JUNK chunk content inspection.
        # Real AVI JUNK chunks are null-padded alignment bytes.
        # Stego-injected JUNK chunks contain printable/structured text.
        # Original code only checked appended_data (outside RIFF); JUNK inside
        # RIFF has matching declared/actual sizes so appended_data returns 0.
        if self.ext == '.avi':
            try:
                raw_avi = open(self.path, 'rb').read()
                # Use find() to locate JUNK chunk anywhere in file (including
                # inside nested LIST chunks which the chunk walker misses)
                junk_pos = raw_avi.find(b'JUNK')
                pos_avi  = junk_pos if junk_pos != -1 else -1
                if pos_avi != -1 and pos_avi + 8 <= len(raw_avi):
                    ctype = b'JUNK'
                    csize = struct.unpack_from('<I', raw_avi, pos_avi + 4)[0]
                    if csize > 64:
                        content   = raw_avi[pos_avi + 8: pos_avi + 8 + min(csize, 2000)]
                        # Genuine padding chunks are all zeros or random; stego is printable
                        printable = sum(1 for b in content if 32 <= b <= 126) / max(len(content), 1)
                        null_ratio = content.count(0) / max(len(content), 1)
                        junk_suspicious = printable > 0.25 and null_ratio < 0.5
                        _junk_preview = (
                            _bytes_text_preview(content, max_chars=1000)
                            if junk_suspicious else {'text': '', 'likely_message': False}
                        )
                        checks['junk_chunk_content'] = {
                            'detected':        junk_suspicious,
                            'score':           60 if junk_suspicious else 0,
                            'chunk_size':      csize,
                            'printable_ratio': round(printable, 3),
                            'null_ratio':      round(null_ratio, 3),
                            'decoded_content': _junk_preview.get('text', ''),
                            'decoded_likely_message': _junk_preview.get('likely_message', False),
                            'info': (
                                f'JUNK chunk size={csize}B, printable={printable:.0%}, nulls={null_ratio:.0%}'
                                f' — {"SUSPICIOUS: structured content in JUNK chunk" if junk_suspicious else "normal padding"}'
                            ),
                        }
                        pass  # done — single JUNK check
            except Exception as e:
                checks['junk_chunk_content'] = {'detected': False, 'score': 0, 'info': f'Error: {e}'}

        # 5. MKV attachments
        if self.ext == '.mkv':
            try:
                raw = open(self.path, 'rb').read(1024 * 1024)  # first 1 MB
                # MKV uses EBML; Attachments element ID is 0x1941A469
                attach_marker = b'\x19\x41\xa4\x69'
                has_attach = attach_marker in raw
                checks['mkv_attachments'] = {
                    'detected': has_attach,
                    'score': 45 if has_attach else 0,
                    'info': 'MKV Attachments element found — may contain embedded files' if has_attach else 'No MKV attachments',
                }
            except Exception as e:
                checks['mkv_attachments'] = {'detected': False, 'score': 0, 'info': f'Error: {e}'}

        return _assemble(self.path, 'Video', checks)


# ─────────────────────────────────────────────────────────────────────────────
#  Public dispatcher
# ─────────────────────────────────────────────────────────────────────────────

def is_image(path):
    return Path(path).suffix.lower() in IMAGE_EXTENSIONS

def is_document(path):
    ext = Path(path).suffix.lower()
    return ext in SUPPORTED_EXTENSIONS and ext not in IMAGE_EXTENSIONS

def analyze_document(path):
    """
    Auto-detect file type and run the correct analyzer.
    Returns a standard dict with keys:
      analyzer_type, file_path, file_type, file_size_bytes, file_hash,
      checks, total_score, risk_level, detected, detected_checks, summary
    """
    ext = Path(path).suffix.lower()

    if ext == '.pdf':
        return PDFAnalyzer(path).run_all()
    elif ext in ('.docx', '.doc', '.odt'):
        return DOCXAnalyzer(path).run_all()
    elif ext in ('.pptx', '.ppt', '.odp'):
        return PPTXAnalyzer(path).run_all()
    elif ext in ('.xlsx', '.xls', '.ods'):
        return XLSXAnalyzer(path).run_all()
    elif ext in PLAIN_TEXT_EXTS:
        return PlainTextAnalyzer(path).run_all()
    elif ext in AUDIO_EXTENSIONS:
        return AudioAnalyzer(path).run_all()
    elif ext in VIDEO_EXTENSIONS:
        return VideoAnalyzer(path).run_all()
    else:
        raise ValueError(
            f'Unsupported extension: {ext}\n'
            f'Supported: {sorted(SUPPORTED_EXTENSIONS)}'
        )
